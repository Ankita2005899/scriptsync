import os
import shutil
import tempfile
import threading
import uuid
import time
import io
import re

import requests
from flask import Flask, render_template, request, jsonify, send_file
from faster_whisper import WhisperModel
from deep_translator import GoogleTranslator
import moviepy.editor as mp

from docx import Document
from docx.shared import Pt, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH

from reportlab.lib.pagesizes import letter
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_LEFT

from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.dml.color import RGBColor as PptxRGBColor

try:
    from duckduckgo_search import DDGS
    print("[startup] duckduckgo_search imported successfully — image search enabled.")
except Exception as e:
    DDGS = None  # image search becomes a no-op if the package fails to load
    print(f"[startup] duckduckgo_search FAILED to import ({type(e).__name__}: {e}) — image search disabled.")

app = Flask(__name__)

# Note: this does NOT change Render free tier's 512MB RAM ceiling.
# A very large video can still crash/hang the worker during transcription.
app.config["MAX_CONTENT_LENGTH"] = 500 * 1024 * 1024  # 500MB

# If set, /api/process forwards the video to a GPU-backed Colab server
# instead of processing it locally on Render's slow free-tier CPU.
# Set this in Render -> Environment -> COLAB_BACKEND_URL to the ngrok URL
# printed by the Colab notebook's last cell (no trailing slash).
# If unset, empty, or unreachable, this app falls back to local CPU
# processing automatically -- nothing breaks if you don't have Colab running.
COLAB_BACKEND_URL = os.environ.get("COLAB_BACKEND_URL", "").rstrip("/")

MODEL_SIZE = os.environ.get("WHISPER_MODEL", "tiny")  # tiny/base = free-tier safe
_model = None
_model_lock = threading.Lock()


def get_model():
    """Load the whisper model once, lazily (keeps cold-start fast)."""
    global _model
    if _model is None:
        with _model_lock:
            if _model is None:
                _model = WhisperModel(MODEL_SIZE, device="cpu", compute_type="int8")
    return _model


def _preload_model():
    """
    Warm up the model as soon as the worker boots, in a background thread,
    instead of on the first user request. Render's free tier wipes any
    cached model files on every restart/redeploy, so this download/load
    cost happens once per deploy no matter what — this just moves it out
    of a user's first request and into server startup instead.
    """
    try:
        get_model()
    except Exception:
        # If preloading fails for any reason, get_model() will just retry
        # lazily on the first real request instead.
        pass


threading.Thread(target=_preload_model, daemon=True).start()


LANGUAGES = {
    "Hindi": "hi", "English": "en", "Marathi": "mr", "Bengali": "bn",
    "Tamil": "ta", "Telugu": "te", "Gujarati": "gu", "Kannada": "kn",
    "Malayalam": "ml", "Punjabi": "pa", "Urdu": "ur",
    "Spanish": "es", "French": "fr", "German": "de", "Portuguese": "pt",
    "Russian": "ru", "Chinese (Simplified)": "zh-CN", "Japanese": "ja",
    "Korean": "ko", "Arabic": "ar",
}

# In-memory job store. Fine for a single free-tier worker;
# jobs are lost on redeploy/restart, which is acceptable here.
JOBS = {}
JOBS_LOCK = threading.Lock()

# Simple cache so we don't re-search the same query twice within one export
_IMAGE_CACHE = {}
_IMAGE_CACHE_LOCK = threading.Lock()

_STOPWORDS = {
    "the", "a", "an", "and", "or", "but", "is", "are", "was", "were", "be",
    "been", "to", "of", "in", "on", "at", "for", "with", "this", "that",
    "it", "as", "by", "from", "so", "we", "you", "i", "he", "she", "they",
    "them", "his", "her", "its", "our", "your", "not", "do", "does", "did",
    "have", "has", "had", "will", "would", "can", "could", "should", "then",
}


def search_image_url(query, timeout=6):
    """
    Looks up one relevant image URL for a short piece of text, using
    DuckDuckGo's free (unofficial, no API key) image search. Returns None
    on any failure -- callers should treat a missing image as fine, not
    fatal, since this library can occasionally be rate-limited or break.
    """
    query = (query or "").strip()
    if not query or DDGS is None:
        print(f"[image_search] SKIPPED — query empty or DDGS unavailable (DDGS={DDGS}), query='{query}'")
        return None

    with _IMAGE_CACHE_LOCK:
        if query in _IMAGE_CACHE:
            return _IMAGE_CACHE[query]

    url = None
    try:
        with DDGS() as ddgs:
            results = list(ddgs.images(query, max_results=1, safesearch="moderate"))
            if results:
                url = results[0].get("image")
                print(f"[image_search] OK — query='{query}' -> {url}")
            else:
                print(f"[image_search] NO RESULTS — query='{query}'")
    except Exception as e:
        print(f"[image_search] EXCEPTION — query='{query}' -> {type(e).__name__}: {e}")
        url = None

    with _IMAGE_CACHE_LOCK:
        _IMAGE_CACHE[query] = url
    return url


def download_image_bytes(url, timeout=8, max_bytes=8 * 1024 * 1024):
    """Downloads an image for embedding into a document. Returns None on failure."""
    if not url:
        return None
    try:
        resp = requests.get(url, timeout=timeout, stream=True)
        if resp.status_code != 200:
            print(f"[image_download] BAD STATUS {resp.status_code} for {url}")
            return None
        content = resp.content
        if not content or len(content) > max_bytes:
            print(f"[image_download] BAD SIZE ({len(content) if content else 0} bytes) for {url}")
            return None
        print(f"[image_download] OK — {len(content)} bytes from {url}")
        return content
    except Exception as e:
        print(f"[image_download] EXCEPTION for {url} -> {type(e).__name__}: {e}")
        return None


def pick_search_query(text, max_words=6):
    """Picks a short, reasonable web-image search query out of a segment's text.
    Unicode-aware so this also works on non-English scripts (Hindi, etc.),
    not just A-Z text -- otherwise translated text in another script
    silently produces an empty query and image search never runs."""
    words = re.findall(r"[^\s.,!?;:()\"'\u2018\u2019\u201c\u201d]+", text, flags=re.UNICODE)
    return " ".join(words[:max_words]) if words else ""


def pick_important_word(text):
    """
    Heuristic 'main word' picker for bolding/highlighting: the longest
    non-stopword in the sentence, on the assumption that longer content
    words tend to carry more meaning than short function words. Uses a
    Unicode-aware pattern so this also works on non-English scripts
    (Hindi, etc.) -- an A-Z-only pattern would find nothing there and
    silently skip highlighting. This is a simple length heuristic, not
    real NLP keyword extraction.
    """
    words = re.findall(r"[^\s.,!?;:()\"'\u2018\u2019\u201c\u201d]+", text, flags=re.UNICODE)
    candidates = [w for w in words if w.lower() not in _STOPWORDS and len(w) > 3]
    if not candidates:
        return None
    return max(candidates, key=len)


def group_segments(segments, char_budget=450, max_per_group=6):
    """
    Groups consecutive segments together so a slide/page shows a
    reasonable amount of content instead of one short segment each.
    Stops a group once it hits char_budget characters of translated
    text or max_per_group segments, whichever comes first.
    """
    groups = []
    current, current_len = [], 0
    for seg in segments:
        seg_len = len(seg.get("translated", ""))
        if current and (current_len + seg_len > char_budget or len(current) >= max_per_group):
            groups.append(current)
            current, current_len = [], 0
        current.append(seg)
        current_len += seg_len
    if current:
        groups.append(current)
    return groups


def make_topic_title(group, max_words=6):
    """
    Derives a short, bold heading-style title from a group of segments --
    e.g. 'Setting Up The Project' -- similar to a slide title in a
    professionally designed deck, instead of just showing a timestamp.
    Unicode-aware so this works for translated text in any script.
    """
    text = " ".join(s.get("translated", "") for s in group)
    words = re.findall(r"[^\s.,!?;:()\"'\u2018\u2019\u201c\u201d]+", text, flags=re.UNICODE)
    if not words:
        return "Untitled Segment"
    title_words = words[:max_words]
    return " ".join(w.capitalize() if w.isascii() else w for w in title_words)


def format_duration(total_seconds):
    """Formats seconds as e.g. '4:32' for stat callouts."""
    m = int(total_seconds // 60)
    s = int(total_seconds % 60)
    return f"{m}:{s:02d}"


def get_job_result_or_404(job_id):
    with JOBS_LOCK:
        job = JOBS.get(job_id)
    if not job or job.get("status") != "done" or not job.get("result"):
        return None
    return job["result"]

# Drop jobs older than this so memory doesn't grow forever (seconds)
JOB_TTL_SECONDS = 60 * 60  # 1 hour

# Ordered stages used to compute a progress percentage on the frontend
STAGE_ORDER = ["queued", "extracting_audio", "loading_model", "transcribing", "translating", "done"]


def _cleanup_old_jobs():
    cutoff = time.time() - JOB_TTL_SECONDS
    with JOBS_LOCK:
        stale = [jid for jid, job in JOBS.items() if job.get("created_at", 0) < cutoff]
        for jid in stale:
            JOBS.pop(jid, None)


def _set_job(job_id, **fields):
    with JOBS_LOCK:
        if job_id in JOBS:
            JOBS[job_id].update(fields)


def _try_colab_backend(job_id, video_path, target_lang):
    """
    Attempts to forward the video to the Colab GPU backend. Returns True if
    it succeeded and the job was completed this way, False if it should
    fall back to local CPU processing instead (Colab not configured, not
    reachable, or it returned an error).
    """
    if not COLAB_BACKEND_URL:
        return False

    _set_job(job_id, status="transcribing")  # Colab is fast enough that finer stages aren't very useful
    try:
        with open(video_path, "rb") as f:
            resp = requests.post(
                f"{COLAB_BACKEND_URL}/api/process",
                files={"video": (os.path.basename(video_path), f)},
                data={"language": target_lang},
                timeout=180,  # Colab GPU should be much faster than this
            )
        if resp.status_code != 200:
            # Colab reachable but returned an error -- fall back to local.
            return False

        data = resp.json()
        if "segments" not in data:
            return False

        _set_job(job_id, status="translating")  # brief, mostly for UI continuity
        _set_job(
            job_id,
            status="done",
            result={
                "detected_language": data.get("detected_language", "unknown"),
                "target_language": data.get("target_language", target_lang),
                "segments": data["segments"],
            },
        )
        return True

    except (requests.exceptions.RequestException, ValueError):
        # Colab notebook not running, tunnel expired, network hiccup, or
        # bad JSON -- silently fall back to local processing below.
        return False


def run_job(job_id, video_path, target_lang, tmp_dir):
    """Runs in a background thread. Does the actual heavy lifting."""
    try:
        if _try_colab_backend(job_id, video_path, target_lang):
            return  # Colab handled it successfully; nothing more to do.

        # Fallback: process locally on Render's CPU, same as before.
        _set_job(job_id, status="extracting_audio")

        audio_path = os.path.join(tmp_dir, "audio.wav")
        clip = mp.VideoFileClip(video_path)
        clip.audio.write_audiofile(audio_path, logger=None)
        clip.close()

        if _model is None:
            _set_job(job_id, status="loading_model")

        model = get_model()

        _set_job(job_id, status="transcribing")
        segments_gen, info = model.transcribe(
            audio_path,
            beam_size=1,
            vad_filter=True,
            condition_on_previous_text=False,
        )
        raw_segments = [
            {"start": round(seg.start, 2), "end": round(seg.end, 2), "original": seg.text.strip()}
            for seg in segments_gen if seg.text.strip()
        ]

        _set_job(job_id, status="translating")

        # Batch translation instead of one API call per segment.
        # deep_translator's GoogleTranslator has a ~5000 char limit per call,
        # so we join segments with a unique delimiter and split the result
        # back apart, cutting dozens of network round-trips down to just a few.
        translator = GoogleTranslator(source="auto", target=target_lang)
        DELIM = "\n|||\n"
        BATCH_CHAR_LIMIT = 4000

        batches = []
        current_batch, current_len = [], 0
        for seg in raw_segments:
            seg_len = len(seg["original"]) + len(DELIM)
            if current_batch and current_len + seg_len > BATCH_CHAR_LIMIT:
                batches.append(current_batch)
                current_batch, current_len = [], 0
            current_batch.append(seg)
            current_len += seg_len
        if current_batch:
            batches.append(current_batch)

        segments = []
        for batch in batches:
            joined = DELIM.join(s["original"] for s in batch)
            try:
                translated_joined = translator.translate(joined)
                parts = translated_joined.split(DELIM.strip())
                if len(parts) != len(batch):
                    # Delimiter got mangled by translation; fall back to
                    # per-segment translation for just this batch.
                    parts = []
                    for s in batch:
                        try:
                            parts.append(translator.translate(s["original"]))
                        except Exception:
                            parts.append(s["original"])
            except Exception:
                parts = [s["original"] for s in batch]

            for seg, translated in zip(batch, parts):
                segments.append({
                    "start": seg["start"],
                    "end": seg["end"],
                    "original": seg["original"],
                    "translated": translated.strip(),
                })

        _set_job(
            job_id,
            status="done",
            result={
                "detected_language": info.language,
                "target_language": target_lang,
                "segments": segments,
            },
        )

    except Exception as e:
        _set_job(job_id, status="error", error=str(e))
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)


def build_docx(result):
    """Clean, professional Word document: a stats-style intro (like a
    report cover), then a bold auto-generated topic heading per group of
    segments, with the translated text as the main body and the original
    as a smaller gray reference line underneath."""
    doc = Document()
    doc.add_heading("ScriptSync Transcript", level=0)

    segments = result.get("segments", [])
    total_duration = segments[-1]["end"] if segments else 0
    groups = group_segments(segments)

    meta = doc.add_paragraph()
    meta_run = meta.add_run(
        f"{result.get('detected_language', 'unknown').upper()} -> "
        f"{result.get('target_language', '').upper()}  \u2022  "
        f"{format_duration(total_duration)} total  \u2022  "
        f"{len(segments)} segments  \u2022  {len(groups)} sections"
    )
    meta_run.italic = True
    meta_run.font.color.rgb = RGBColor(0x80, 0x80, 0x80)

    for i, group in enumerate(groups, start=1):
        doc.add_heading(f"{i}. {make_topic_title(group)}", level=1)

        time_p = doc.add_paragraph()
        time_run = time_p.add_run(f"{group[0]['start']:.0f}s \u2013 {group[-1]['end']:.0f}s")
        time_run.font.size = Pt(9)
        time_run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
        time_run.italic = True

        # Merge every segment's translated text into ONE flowing paragraph
        # (like real prose in a document) instead of a separate short line
        # per segment -- each segment's key word is still emphasized inline.
        body_p = doc.add_paragraph()
        for j, seg in enumerate(group):
            translated = seg.get("translated", "")
            important = pick_important_word(translated)
            _add_emphasized_docx_run(body_p, translated, important, base_size=13, emphasis_size=15)
            if j < len(group) - 1:
                body_p.add_run(" ").font.size = Pt(13)

        original_combined = " ".join(seg.get("original", "") for seg in group)
        original_p = doc.add_paragraph()
        orig_run = original_p.add_run(f"Original: {original_combined}")
        orig_run.font.size = Pt(9)
        orig_run.font.color.rgb = RGBColor(0x80, 0x80, 0x80)
        orig_run.italic = True

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf


def build_pdf(result):
    """PDF export with the main (heuristically important) word of each
    segment bolded, and one auto-searched image per group of segments."""
    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=letter, topMargin=0.75 * inch, bottomMargin=0.75 * inch)
    styles = getSampleStyleSheet()

    title_style = ParagraphStyle("TitleStyle", parent=styles["Title"], alignment=TA_LEFT)
    body_style = ParagraphStyle("BodyStyle", parent=styles["Normal"], fontSize=12, leading=17, spaceAfter=4)
    meta_style = ParagraphStyle("MetaStyle", parent=styles["Normal"], fontSize=9, textColor="#666666", spaceAfter=14)
    time_style = ParagraphStyle("TimeStyle", parent=styles["Normal"], fontSize=9, textColor="#B3061B", spaceBefore=10)

    story = [Paragraph("ScriptSync Transcript", title_style), Spacer(1, 6)]
    story.append(Paragraph(
        f"Detected: {result.get('detected_language', 'unknown').upper()} &rarr; "
        f"Translated to: {result.get('target_language', '').upper()}",
        meta_style,
    ))

    from reportlab.platypus import Image as RLImage

    for group in group_segments(result.get("segments", [])):
        for seg in group:
            story.append(Paragraph(f"{seg['start']:.0f}s – {seg['end']:.0f}s", time_style))

            translated = seg.get("translated", "")
            important = pick_important_word(translated)
            display_text = translated
            if important:
                # Bold only the first occurrence of the chosen important word.
                display_text = re.sub(
                    rf"\b({re.escape(important)})\b",
                    r"<b>\1</b>",
                    translated,
                    count=1,
                )
            story.append(Paragraph(display_text, body_style))

        # One image per group, not per segment, so pages don't fill up
        # with mostly-whitespace single-line entries.
        query = pick_search_query(" ".join(s.get("original", "") for s in group))
        img_url = search_image_url(query)
        img_bytes = download_image_bytes(img_url)
        if img_bytes:
            try:
                story.append(Spacer(1, 6))
                story.append(RLImage(io.BytesIO(img_bytes), width=2.8 * inch, height=1.8 * inch))
                story.append(Spacer(1, 10))
            except Exception:
                pass  # bad/corrupt image data -- just skip it for this group

    doc.build(story)
    buf.seek(0)
    return buf


def _add_emphasized_run(paragraph, text, keyword, base_size, emphasis_size):
    """
    Adds `text` to a pptx paragraph as normal-size runs, except the first
    occurrence of `keyword` (if any) which is bold and shown larger --
    so only the important word stands out, not the whole sentence.
    """
    if not keyword:
        run = paragraph.add_run()
        run.text = text
        run.font.size = PptxPt(base_size)
        return

    match = re.search(rf"\b{re.escape(keyword)}\b", text)
    if not match:
        run = paragraph.add_run()
        run.text = text
        run.font.size = PptxPt(base_size)
        return

    before, kw, after = text[:match.start()], match.group(0), text[match.end():]
    if before:
        r = paragraph.add_run()
        r.text = before
        r.font.size = PptxPt(base_size)
    r_kw = paragraph.add_run()
    r_kw.text = kw
    r_kw.font.size = PptxPt(emphasis_size)
    r_kw.font.bold = True
    r_kw.font.color.rgb = PptxRGBColor(0xD6, 0x27, 0x3C)
    if after:
        r = paragraph.add_run()
        r.text = after
        r.font.size = PptxPt(base_size)


def build_pptx(result, with_images=True):
    """Professional-style deck: a title slide with stat callouts (like the
    reference design), then content slides with a bold auto-generated
    topic heading, subtitle-style time range, grouped segment text with
    the key word emphasized, and a prominent auto-searched image."""
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]

    segments = result.get("segments", [])
    total_duration = segments[-1]["end"] if segments else 0
    groups = group_segments(segments, char_budget=380, max_per_group=5)

    # ---- Title slide, styled with stat callouts like the reference deck ----
    title_slide = prs.slides.add_slide(blank_layout)

    title_box = title_slide.shapes.add_textbox(PptxInches(0.6), PptxInches(0.7), PptxInches(9), PptxInches(1.2))
    title_run = title_box.text_frame.paragraphs[0].add_run()
    title_run.text = "ScriptSync Transcript"
    title_run.font.size = PptxPt(40)
    title_run.font.bold = True
    title_run.font.color.rgb = PptxRGBColor(0x1A, 0x0E, 0x10)

    subtitle_box = title_slide.shapes.add_textbox(PptxInches(0.6), PptxInches(1.55), PptxInches(9), PptxInches(0.6))
    subtitle_run = subtitle_box.text_frame.paragraphs[0].add_run()
    subtitle_run.text = (
        f"{result.get('detected_language', 'unknown').upper()} -> "
        f"{result.get('target_language', '').upper()}"
    )
    subtitle_run.font.size = PptxPt(16)
    subtitle_run.font.color.rgb = PptxRGBColor(0x80, 0x80, 0x80)

    stats = [
        (format_duration(total_duration), "Total duration"),
        (str(len(segments)), "Segments transcribed"),
        (str(len(groups)), "Slides generated"),
    ]
    stat_x = 0.6
    for value, label in stats:
        box = title_slide.shapes.add_textbox(PptxInches(stat_x), PptxInches(2.6), PptxInches(2.8), PptxInches(1.2))
        tf = box.text_frame
        v_run = tf.paragraphs[0].add_run()
        v_run.text = value
        v_run.font.size = PptxPt(30)
        v_run.font.bold = True
        v_run.font.color.rgb = PptxRGBColor(0xD6, 0x27, 0x3C)

        l_p = tf.add_paragraph()
        l_run = l_p.add_run()
        l_run.text = label
        l_run.font.size = PptxPt(11)
        l_run.font.color.rgb = PptxRGBColor(0x80, 0x80, 0x80)
        stat_x += 3.0

    # ---- Content slides ----
    for group in groups:
        slide = prs.slides.add_slide(blank_layout)

        heading_text = make_topic_title(group)
        time_range = f"{group[0]['start']:.0f}s – {group[-1]['end']:.0f}s"

        text_width = PptxInches(5.7) if with_images else PptxInches(9)

        heading_box = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.4), text_width, PptxInches(0.7))
        h_run = heading_box.text_frame.paragraphs[0].add_run()
        h_run.text = heading_text
        h_run.font.size = PptxPt(26)
        h_run.font.bold = True
        h_run.font.color.rgb = PptxRGBColor(0x1A, 0x0E, 0x10)

        sub_box = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(1.0), text_width, PptxInches(0.4))
        sub_run = sub_box.text_frame.paragraphs[0].add_run()
        sub_run.text = time_range
        sub_run.font.size = PptxPt(12)
        sub_run.font.color.rgb = PptxRGBColor(0x80, 0x80, 0x80)

        tb = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(1.5), text_width, PptxInches(5.4))
        tf = tb.text_frame
        tf.word_wrap = True

        # Merge every segment's translated text into ONE flowing paragraph
        # (like real slide body copy) instead of choppy separate lines --
        # each segment's key word is still emphasized inline.
        body_p = tf.paragraphs[0]
        for j, seg in enumerate(group):
            translated = seg.get("translated", "")
            important = pick_important_word(translated)
            _add_emphasized_run(body_p, translated, important, base_size=15, emphasis_size=18)
            if j < len(group) - 1:
                sp = body_p.add_run()
                sp.text = " "
                sp.font.size = PptxPt(15)

        original_combined = " ".join(seg.get("original", "") for seg in group)
        orig_p = tf.add_paragraph()
        orig_run = orig_p.add_run()
        orig_run.text = original_combined
        orig_run.font.size = PptxPt(10)
        orig_run.font.italic = True
        orig_run.font.color.rgb = PptxRGBColor(0x99, 0x99, 0x99)

        if with_images:
            query = pick_search_query(" ".join(s.get("original", "") for s in group))
            img_url = search_image_url(query)
            img_bytes = download_image_bytes(img_url)
            if img_bytes:
                try:
                    slide.shapes.add_picture(
                        io.BytesIO(img_bytes),
                        PptxInches(6.4), PptxInches(1.5),
                        width=PptxInches(3.0),
                    )
                except Exception:
                    pass  # bad/corrupt image data -- just skip it for this slide

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def _add_emphasized_docx_run(paragraph, text, keyword, base_size, emphasis_size):
    """
    Same idea as the pptx version: adds `text` to a docx paragraph as
    normal runs, except the first occurrence of `keyword` which is bold,
    highlighted, and shown larger -- so only the key word stands out.
    """
    from docx.enum.text import WD_COLOR_INDEX

    if not keyword:
        r = paragraph.add_run(text)
        r.font.size = Pt(base_size)
        return

    match = re.search(rf"\b{re.escape(keyword)}\b", text)
    if not match:
        r = paragraph.add_run(text)
        r.font.size = Pt(base_size)
        return

    before, kw, after = text[:match.start()], match.group(0), text[match.end():]
    if before:
        r = paragraph.add_run(before)
        r.font.size = Pt(base_size)
    r_kw = paragraph.add_run(kw)
    r_kw.font.size = Pt(emphasis_size)
    r_kw.bold = True
    try:
        r_kw.font.highlight_color = WD_COLOR_INDEX.YELLOW
    except Exception:
        pass
    if after:
        r = paragraph.add_run(after)
        r.font.size = Pt(base_size)


def build_notes_docx(result):
    """
    'Notes' format: groups several segments per block (instead of one
    segment each, which left too much empty space), highlights only the
    key word per segment rather than the whole sentence, and adds one
    auto-searched image per group.
    """
    doc = Document()
    doc.add_heading("ScriptSync — Notes", level=0)

    meta = doc.add_paragraph()
    meta.add_run(
        f"{result.get('detected_language', 'unknown').upper()} -> "
        f"{result.get('target_language', '').upper()}"
    ).italic = True

    for group in group_segments(result.get("segments", []), char_budget=500, max_per_group=6):
        for seg in group:
            time_p = doc.add_paragraph()
            time_run = time_p.add_run(f"{seg['start']:.0f}s – {seg['end']:.0f}s")
            time_run.font.size = Pt(9)
            time_run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
            time_run.bold = True

            translated = seg.get("translated", "")
            important = pick_important_word(translated)
            note_p = doc.add_paragraph()
            _add_emphasized_docx_run(note_p, translated, important, base_size=12, emphasis_size=15)

        # One image per group of segments, not one per single short
        # segment, so the page isn't mostly whitespace around tiny blocks.
        query = pick_search_query(" ".join(s.get("original", "") for s in group))
        img_url = search_image_url(query)
        img_bytes = download_image_bytes(img_url)
        if img_bytes:
            try:
                doc.add_picture(io.BytesIO(img_bytes), width=Inches(2.8))
            except Exception:
                pass

        doc.add_paragraph()  # spacing between groups

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf


@app.route("/")
def index():
    return render_template("index.html", languages=LANGUAGES)


@app.route("/api/ready")
def api_ready():
    return jsonify({"model_ready": _model is not None})


@app.route("/api/languages")
def api_languages():
    return jsonify(LANGUAGES)


@app.route("/api/process", methods=["POST"])
def api_process():
    _cleanup_old_jobs()

    if "video" not in request.files:
        return jsonify({"error": "No video file received."}), 400

    video = request.files["video"]
    target_lang = request.form.get("language", "hi")

    if video.filename == "":
        return jsonify({"error": "Empty filename."}), 400

    tmp_dir = tempfile.mkdtemp()
    video_path = os.path.join(tmp_dir, video.filename)
    video.save(video_path)

    job_id = str(uuid.uuid4())
    with JOBS_LOCK:
        JOBS[job_id] = {"status": "queued", "created_at": time.time()}

    thread = threading.Thread(
        target=run_job,
        args=(job_id, video_path, target_lang, tmp_dir),
        daemon=True,
    )
    thread.start()

    # Return immediately so Render's proxy never has to hold this request open.
    return jsonify({"job_id": job_id}), 202


@app.route("/api/export/<job_id>/<fmt>")
def api_export(job_id, fmt):
    result = get_job_result_or_404(job_id)
    if result is None:
        return jsonify({"error": "This job isn't finished yet, or has expired. Please generate the script again."}), 404

    lang = result.get("target_language", "output")

    try:
        if fmt == "docx":
            buf = build_docx(result)
            return send_file(buf, as_attachment=True, download_name=f"scriptsync_{lang}.docx",
                              mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document")

        if fmt == "pdf":
            buf = build_pdf(result)
            return send_file(buf, as_attachment=True, download_name=f"scriptsync_{lang}.pdf",
                              mimetype="application/pdf")

        if fmt == "pptx":
            buf = build_pptx(result, with_images=True)
            return send_file(buf, as_attachment=True, download_name=f"scriptsync_{lang}.pptx",
                              mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation")

        if fmt == "notes":
            buf = build_notes_docx(result)
            return send_file(buf, as_attachment=True, download_name=f"scriptsync_notes_{lang}.docx",
                              mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document")

        return jsonify({"error": f"Unknown format '{fmt}'."}), 400

    except Exception as e:
        return jsonify({"error": f"Failed to generate {fmt}: {e}"}), 500


@app.route("/api/status/<job_id>")
def api_status(job_id):
    with JOBS_LOCK:
        job = JOBS.get(job_id)

    if job is None:
        return jsonify({"error": "Unknown or expired job."}), 404

    status = job.get("status")
    try:
        stage_index = STAGE_ORDER.index(status)
    except ValueError:
        stage_index = 0
    progress_pct = round((stage_index / (len(STAGE_ORDER) - 1)) * 100)

    response = {"status": status, "progress": progress_pct}
    if status == "done":
        response["result"] = job.get("result")
    elif status == "error":
        response["error"] = job.get("error")

    return jsonify(response)


if __name__ == "__main__":
    app.run(host="0.0.0.0", port=int(os.environ.get("PORT", 5000)))
